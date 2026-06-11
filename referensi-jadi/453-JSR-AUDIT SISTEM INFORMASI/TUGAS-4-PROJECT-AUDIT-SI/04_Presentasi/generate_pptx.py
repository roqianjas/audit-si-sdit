"""
Generate Slide Presentasi Project Audit SI — PT. Murni Solusindo Nusantara
Framework COBIT 2019 | 17 Slide | Dark Navy Theme
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os

# === PATHS ===
BASE = r"c:\Users\adminlocal\Documents\UNM\UNM-Semester-8\453-JSR-AUDIT SISTEM INFORMASI\TUGAS-4-PROJECT-AUDIT-SI"
OUTPUT = os.path.join(BASE, "04_Presentasi", "Kelompok-Presentasi-Audit-SI-PT_Murni.pptx")

# === COLORS ===
C_ACCENT  = RGBColor(0x00, 0x96, 0xD6)   # Bright blue accent
C_ACCENT2 = RGBColor(0x00, 0xBC, 0x8C)   # Teal green
C_BG      = RGBColor(0x0D, 0x1B, 0x2A)   # Dark navy
C_BG2     = RGBColor(0x1B, 0x2A, 0x41)   # Slightly lighter navy
C_BG3     = RGBColor(0x15, 0x3B, 0x60)   # Card blue
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xBB, 0xBB, 0xBB)
C_ORANGE  = RGBColor(0xFF, 0x8C, 0x00)
C_RED     = RGBColor(0xE5, 0x39, 0x35)
C_YELLOW  = RGBColor(0xFF, 0xD6, 0x00)
C_GREEN   = RGBColor(0x2E, 0x7D, 0x32)
C_GOLD    = RGBColor(0xC6, 0x8A, 0x2C)

TOTAL = 17

# === INIT ===
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# === HELPERS ===
def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    # Left accent bar
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    r.fill.solid(); r.fill.fore_color.rgb = C_ACCENT; r.line.fill.background()
    return s

def box(slide, l, t, w, h, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()
    return r

def rounded_box(slide, l, t, w, h, color):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()
    return r

def txt(slide, l, t, w, h, text, sz=18, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold; p.font.name = 'Calibri'
    p.alignment = align
    return tb

def multi_txt(slide, l, t, w, h, lines, sz=14, color=C_WHITE, bold=False, spacing=6):
    """Add a textbox with multiple paragraphs."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold
        p.font.name = 'Calibri'; p.space_after = Pt(spacing)
    return tb

def bullets(slide, l, t, w, h, items, sz=14, color=C_WHITE, spacing=8):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = 'Calibri'; p.space_after = Pt(spacing)

def snum(slide, n):
    txt(slide, 12.3, 7.0, 1, 0.4, f"{n}/{TOTAL}", sz=10, color=C_LIGHT, align=PP_ALIGN.RIGHT)

def aline(slide, l, t, w):
    box(slide, l, t, w, 0.04, C_ACCENT)

def header(slide, label, title):
    txt(slide, 1, 0.4, 5, 0.4, label, sz=14, color=C_ACCENT, bold=True)
    txt(slide, 1, 0.85, 11, 0.6, title, sz=28, color=C_WHITE, bold=True)
    aline(slide, 1, 1.45, 3)


# ============================================================
# SLIDE 1: COVER
# ============================================================
s = new_slide()
# Top stripe
box(s, 0, 0, 13.333, 0.6, C_ACCENT)
txt(s, 1, 0.1, 11, 0.4, "LAPORAN PROJECT AUDIT SISTEM INFORMASI", sz=14, color=C_BG, bold=True, align=PP_ALIGN.CENTER)

txt(s, 1, 1.2, 11, 0.5, "MATA KULIAH AUDIT SISTEM INFORMASI (453)", sz=14, color=C_ACCENT, bold=True)
aline(s, 1, 1.7, 2.5)

txt(s, 1, 2.0, 11, 1.5, "Evaluasi Tata Kelola Teknologi Informasi\nMenggunakan Framework COBIT 2019", sz=36, color=C_WHITE, bold=True)
txt(s, 1, 3.6, 11, 0.6, "Pada PT. Murni Solusindo Nusantara", sz=22, color=C_LIGHT)

txt(s, 1, 4.8, 5, 0.4, "Disusun Oleh:", sz=14, color=C_ACCENT, bold=True)
for i, m in enumerate(["Roki Anjas (11250066)", "Susanto (11250068)"]):
    txt(s, 1, 5.2 + i*0.35, 5, 0.3, m, sz=13, color=C_LIGHT)

txt(s, 1, 6.3, 5, 0.4, "Dosen Pengampu: Juarni Siregar, S.Pd., M.Kom", sz=11, color=C_LIGHT)
txt(s, 1, 6.7, 5, 0.3, "Program Studi Sistem Informasi (S1) | Universitas Nusa Mandiri | 2026", sz=11, color=C_LIGHT)

# Right: Company icon
box(s, 9, 2.0, 3.8, 3.8, C_BG2)
txt(s, 9, 2.7, 3.8, 1.5, "🏢", sz=72, color=C_ACCENT, align=PP_ALIGN.CENTER)
txt(s, 9, 4.3, 3.8, 0.8, "PT. MURNI\nSOLUSINDO\nNUSANTARA", sz=14, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
snum(s, 1)


# ============================================================
# SLIDE 2: AGENDA
# ============================================================
s = new_slide()
txt(s, 1, 0.8, 11, 0.8, "AGENDA PRESENTASI", sz=32, color=C_WHITE, bold=True)
aline(s, 1, 1.6, 3)
outline = [
    ("01", "Pendahuluan",          "Latar belakang, permasalahan, dan tujuan audit"),
    ("02", "Profil Perusahaan",    "PT. Murni Solusindo Nusantara"),
    ("03", "Framework COBIT 2019", "Penjelasan framework dan domain yang dipilih"),
    ("04", "Metodologi",           "RACI Matrix dan kuesioner penelitian"),
    ("05", "Hasil Audit",          "Capability level dan GAP analysis"),
    ("06", "Rekomendasi",          "Perbaikan untuk meningkatkan capability level"),
    ("07", "Kesimpulan & Saran",   "Ringkasan temuan dan saran"),
]
for idx, (num, title, desc) in enumerate(outline):
    y = 2.2 + idx * 0.7
    box(s, 1, y, 0.6, 0.5, C_ACCENT)
    txt(s, 1, y+0.05, 0.6, 0.4, num, sz=14, color=C_BG, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 1.8, y, 4, 0.3, title, sz=18, color=C_WHITE, bold=True)
    txt(s, 1.8, y+0.3, 8, 0.3, desc, sz=12, color=C_LIGHT)
snum(s, 2)


# ============================================================
# SLIDE 3: LATAR BELAKANG
# ============================================================
s = new_slide()
header(s, "01 | PENDAHULUAN", "Latar Belakang")
txt(s, 1, 1.8, 11, 1.2,
    "PT. Murni Solusindo Nusantara adalah perusahaan distributor alat hitung uang dan alat deteksi "
    "uang palsu yang juga memiliki produk digital: Platform Web Builder SaaS, Platform Listing Property, "
    "dan Platform Indoconnex. Dengan portofolio digital yang terus berkembang, diperlukan audit tata kelola TI "
    "untuk memastikan keamanan, efisiensi, dan efektivitas pengelolaan sistem informasi.",
    sz=14, color=C_LIGHT)

for i, (icon, t, d) in enumerate([
    ("⚠️", "Risiko Keamanan Data", "Platform digital menyimpan data pelanggan yang sensitif"),
    ("🔄", "Manajemen Perubahan", "Perubahan sistem web development perlu tata kelola yang terstruktur"),
    ("📊", "Monitoring Operasional", "Operasional platform SaaS harus stabil dan termonitor"),
    ("🛡️", "Standarisasi Prosedur", "SOP tata kelola TI belum terdokumentasi secara formal"),
]):
    y = 3.5 + i * 0.95
    box(s, 1, y, 11.5, 0.8, C_BG2)
    txt(s, 1.3, y+0.05, 0.5, 0.35, icon, sz=20, color=C_WHITE)
    txt(s, 1.9, y+0.05, 5, 0.35, t, sz=14, color=C_ACCENT, bold=True)
    txt(s, 1.9, y+0.4, 10, 0.35, d, sz=12, color=C_LIGHT)
snum(s, 3)


# ============================================================
# SLIDE 4: RUMUSAN MASALAH & TUJUAN
# ============================================================
s = new_slide()
header(s, "01 | PENDAHULUAN", "Rumusan Masalah & Tujuan")

# Left: Rumusan masalah
txt(s, 1, 1.8, 5.5, 0.4, "RUMUSAN MASALAH", sz=18, color=C_RED, bold=True)
aline(s, 1, 2.2, 2)
bullets(s, 1, 2.5, 5.5, 3, [
    "Bagaimana capability level tata kelola TI pada 5 domain COBIT 2019?",
    "Seberapa besar GAP antara kondisi as-is dengan to-be?",
    "Rekomendasi perbaikan apa yang dapat diberikan?",
], sz=13, color=C_LIGHT, spacing=14)

# Right: Tujuan
txt(s, 7, 1.8, 5.5, 0.4, "TUJUAN AUDIT", sz=18, color=C_ACCENT2, bold=True)
aline(s, 7, 2.2, 2)
bullets(s, 7, 2.5, 5.5, 3, [
    "Mengukur capability level pada 5 domain COBIT 2019",
    "Mengidentifikasi GAP antara as-is vs to-be",
    "Menyusun rekomendasi perbaikan yang implementatif",
], sz=13, color=C_LIGHT, spacing=14)
snum(s, 4)


# ============================================================
# SLIDE 5: PROFIL PERUSAHAAN
# ============================================================
s = new_slide()
header(s, "02 | PROFIL PERUSAHAAN", "PT. Murni Solusindo Nusantara")

for i, (lbl, val) in enumerate([
    ("🏢  Nama",        "PT. Murni Solusindo Nusantara"),
    ("📋  Bidang Usaha", "Distributor alat hitung uang & deteksi uang palsu"),
    ("💻  Produk Digital","Web Builder SaaS, Listing Property, Indoconnex"),
    ("👥  Departemen",   "HR, Finance, IT, Digital Marketing"),
    ("🎯  Fokus Audit",  "Divisi IT & Digital Marketing (Web Dev & SEO)"),
]):
    y = 1.9 + i*0.75
    box(s, 1, y, 5.5, 0.6, C_BG2)
    txt(s, 1.3, y+0.1, 2.5, 0.4, lbl, sz=13, color=C_ACCENT, bold=True)
    txt(s, 3.8, y+0.1, 3, 0.4, val, sz=13, color=C_WHITE)

# Right: Responden
box(s, 7, 1.9, 5.5, 5.1, C_BG2)
txt(s, 7, 2.0, 5.5, 0.5, "RESPONDEN KUESIONER", sz=16, color=C_ACCENT, bold=True, align=PP_ALIGN.CENTER)
responden = [
    ("Kepala Divisi IT", "1 orang"),
    ("Senior Web Developer", "1 orang"),
    ("Junior Web Developer", "1 orang"),
    ("Content & SEO Specialist", "1 orang"),
    ("Staff IT Support", "1 orang"),
    ("Digital Marketing Specialist", "1 orang"),
    ("Staff Finance", "1 orang"),
    ("Staff HR", "1 orang"),
]
for i, (jab, jml) in enumerate(responden):
    y = 2.6 + i*0.40
    txt(s, 7.3, y, 3.5, 0.3, f"•  {jab}", sz=11, color=C_LIGHT)
    txt(s, 10.8, y, 1.5, 0.3, jml, sz=11, color=C_ACCENT, bold=True, align=PP_ALIGN.RIGHT)

txt(s, 7, 6.1, 5.5, 0.5, "Total: 8 Responden", sz=14, color=C_YELLOW, bold=True, align=PP_ALIGN.CENTER)
snum(s, 5)


# ============================================================
# SLIDE 6: FRAMEWORK COBIT 2019
# ============================================================
s = new_slide()
header(s, "03 | FRAMEWORK", "COBIT 2019 — Overview")

txt(s, 1, 1.8, 11, 0.6,
    "COBIT 2019 dikembangkan oleh ISACA sebagai framework tata kelola dan manajemen TI terkini. "
    "Berbasis CMMI, fleksibel, dan mengintegrasikan ITIL, TOGAF, serta standar internasional lainnya.",
    sz=14, color=C_LIGHT)

for i, (num, title, desc) in enumerate([
    ("40", "Proses", "5 Governance + 35 Management"),
    ("6", "Prinsip", "Sistem & Kerangka Tata Kelola"),
    ("0-5", "Cap. Level", "Berbasis CMMI"),
    ("N/P/L/F", "Rating", "Not/Partially/Largely/Fully Achieved"),
]):
    l = 1 + i * 3.05
    box(s, l, 2.8, 2.8, 2.0, C_BG2)
    txt(s, l, 2.9, 2.8, 0.7, num, sz=32, color=C_ACCENT, bold=True, align=PP_ALIGN.CENTER)
    txt(s, l, 3.6, 2.8, 0.4, title, sz=14, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, l, 4.0, 2.8, 0.6, desc, sz=11, color=C_LIGHT, align=PP_ALIGN.CENTER)

# Capability Level table
txt(s, 1, 5.2, 4, 0.4, "CAPABILITY LEVEL", sz=16, color=C_ACCENT, bold=True)
levels = [
    ("0", "Incomplete", C_RED),
    ("1", "Performed", C_ORANGE),
    ("2", "Managed", C_YELLOW),
    ("3", "Established", C_ACCENT2),
    ("4", "Predictable", C_ACCENT),
    ("5", "Optimizing", C_GREEN),
]
for i, (lv, name, clr) in enumerate(levels):
    l = 1 + i * 2.0
    box(s, l, 5.7, 1.8, 0.5, clr)
    txt(s, l, 5.72, 0.6, 0.4, lv, sz=14, color=C_BG, bold=True, align=PP_ALIGN.CENTER)
    txt(s, l + 0.5, 5.72, 1.3, 0.4, name, sz=11, color=C_BG, bold=True)
snum(s, 6)


# ============================================================
# SLIDE 7: DOMAIN YANG DIPILIH
# ============================================================
s = new_slide()
header(s, "03 | FRAMEWORK", "Domain COBIT 2019 yang Diaudit")

domains = [
    ("APO12", "Managed Risk", "Manajemen risiko platform digital", "⚠️"),
    ("APO13", "Managed Security", "Keamanan informasi & data pelanggan", "🔒"),
    ("BAI06", "Managed IT Changes", "Change management web development", "🔄"),
    ("DSS01", "Managed Operations", "Operasional harian platform SaaS", "⚙️"),
    ("DSS05", "Managed Security Services", "Layanan keamanan infrastruktur digital", "🛡️"),
]
for i, (code, name, desc, icon) in enumerate(domains):
    y = 1.9 + i * 1.05
    box(s, 1, y, 11.5, 0.9, C_BG2)
    # Icon circle
    box(s, 1.2, y+0.1, 0.7, 0.7, C_BG3)
    txt(s, 1.2, y+0.12, 0.7, 0.6, icon, sz=22, color=C_WHITE, align=PP_ALIGN.CENTER)
    # Code badge
    box(s, 2.1, y+0.15, 1.2, 0.35, C_ACCENT)
    txt(s, 2.1, y+0.17, 1.2, 0.3, code, sz=12, color=C_BG, bold=True, align=PP_ALIGN.CENTER)
    # Name & description
    txt(s, 3.5, y+0.1, 4, 0.35, name, sz=16, color=C_WHITE, bold=True)
    txt(s, 3.5, y+0.45, 9, 0.35, desc, sz=12, color=C_LIGHT)
snum(s, 7)


# ============================================================
# SLIDE 8: RACI MATRIX
# ============================================================
s = new_slide()
header(s, "04 | METODOLOGI", "RACI Matrix")

txt(s, 1, 1.8, 11, 0.5,
    "RACI Matrix mendefinisikan peran dan tanggung jawab setiap pihak dalam pelaksanaan audit.",
    sz=14, color=C_LIGHT)

# RACI explanation
for i, (letter, full, desc, clr) in enumerate([
    ("R", "Responsible", "Melaksanakan tugas", C_ACCENT),
    ("A", "Accountable", "Bertanggung jawab atas keputusan", C_ACCENT2),
    ("C", "Consulted", "Dimintai pendapat", C_YELLOW),
    ("I", "Informed", "Diberikan informasi", C_LIGHT),
]):
    l = 1 + i * 3.05
    box(s, l, 2.6, 0.5, 0.5, clr)
    txt(s, l, 2.62, 0.5, 0.4, letter, sz=18, color=C_BG, bold=True, align=PP_ALIGN.CENTER)
    txt(s, l+0.6, 2.6, 2, 0.25, full, sz=13, color=C_WHITE, bold=True)
    txt(s, l+0.6, 2.85, 2, 0.25, desc, sz=10, color=C_LIGHT)

# RACI table content
raci_data = [
    ("APO12 — Managed Risk",     "A", "R", "R", "C", "I", "I"),
    ("APO13 — Managed Security",  "A", "R", "R", "I", "I", "C"),
    ("BAI06 — Managed IT Changes","A", "R", "R", "I", "I", "I"),
    ("DSS01 — Managed Operations","A", "R", "R", "I", "I", "I"),
    ("DSS05 — Managed Security",  "A", "R", "R", "I", "C", "I"),
]
roles = ["Ka. IT", "Web Dev", "IT Supp", "Dig.Mkt", "Finance", "HR"]

# Header row
cols_start = 4.5
col_w = 1.3
box(s, 1, 3.5, 3.3, 0.45, C_ACCENT)
txt(s, 1.1, 3.52, 3.1, 0.4, "Domain / Aktivitas", sz=11, color=C_BG, bold=True)
for j, role in enumerate(roles):
    box(s, cols_start + j*col_w, 3.5, col_w-0.05, 0.45, C_ACCENT)
    txt(s, cols_start + j*col_w, 3.52, col_w-0.05, 0.4, role, sz=10, color=C_BG, bold=True, align=PP_ALIGN.CENTER)

# Data rows
for i, row in enumerate(raci_data):
    y = 4.0 + i * 0.55
    bg = C_BG2 if i % 2 == 0 else C_BG
    box(s, 1, y, 3.3, 0.5, bg)
    txt(s, 1.1, y+0.05, 3.1, 0.4, row[0], sz=10, color=C_LIGHT)
    for j, val in enumerate(row[1:]):
        box(s, cols_start + j*col_w, y, col_w-0.05, 0.5, bg)
        clr_map = {"R": C_ACCENT, "A": C_ACCENT2, "C": C_YELLOW, "I": C_LIGHT}
        txt(s, cols_start + j*col_w, y+0.05, col_w-0.05, 0.4, val,
            sz=14, color=clr_map.get(val, C_LIGHT), bold=True, align=PP_ALIGN.CENTER)
snum(s, 8)


# ============================================================
# SLIDE 9: METODOLOGI PERHITUNGAN
# ============================================================
s = new_slide()
header(s, "04 | METODOLOGI", "Alur Perhitungan Capability Level")

steps = [
    ("1", "KUESIONER", "Skala Likert 1-5\n(STS, TS, R, S, SS)", C_ACCENT),
    ("2", "HITUNG\nSKOR", "Total skor per\npertanyaan", C_BG3),
    ("3", "HITUNG\nMEAN", "Mean = Total /\nJumlah Responden", C_BG3),
    ("4", "% CAPAI", "% = (Mean/5)\n× 100%", C_BG3),
    ("5", "RATING", "N/P/L/F\nberdasarkan %", C_ORANGE),
    ("6", "CAP.\nLEVEL", "Level tertinggi\ndengan rating L/F", C_ACCENT2),
    ("7", "GAP", "To-be − As-is\n= Selisih Level", C_RED),
]
for i, (num, title, desc, clr) in enumerate(steps):
    l = 0.7 + i * 1.78
    box(s, l, 2.0, 1.6, 2.5, clr)
    txt(s, l, 2.1, 1.6, 0.5, num, sz=28, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, l, 2.6, 1.6, 0.6, title, sz=12, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, l, 3.3, 1.6, 0.8, desc, sz=10, color=C_LIGHT, align=PP_ALIGN.CENTER)
    if i < 6:
        txt(s, l+1.55, 2.8, 0.3, 0.5, "→", sz=20, color=C_ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Rating scale
txt(s, 1, 5.0, 11, 0.4, "RATING SCALE", sz=16, color=C_ACCENT, bold=True)
for i, (code, name, pct, clr) in enumerate([
    ("N", "Not Achieved", "0–15%", C_RED),
    ("P", "Partially Achieved", ">15–50%", C_ORANGE),
    ("L", "Largely Achieved", ">50–85%", C_YELLOW),
    ("F", "Fully Achieved", ">85–100%", C_GREEN),
]):
    l = 1 + i * 3.0
    box(s, l, 5.5, 0.5, 0.5, clr)
    txt(s, l, 5.52, 0.5, 0.4, code, sz=18, color=C_BG, bold=True, align=PP_ALIGN.CENTER)
    txt(s, l+0.6, 5.5, 2, 0.25, name, sz=12, color=C_WHITE, bold=True)
    txt(s, l+0.6, 5.75, 2, 0.25, pct, sz=11, color=C_LIGHT)
snum(s, 9)


# ============================================================
# SLIDE 10: HASIL — CAPABILITY LEVEL PER DOMAIN
# ============================================================
s = new_slide()
header(s, "05 | HASIL AUDIT", "Capability Level Per Domain")

results = [
    ("APO12", "Managed Risk", 60.83, 51.67, 45.00, 2),
    ("APO13", "Managed Security", 65.00, 54.17, 46.67, 2),
    ("BAI06", "Managed IT Changes", 62.50, 53.33, 46.67, 2),
    ("DSS01", "Managed Operations", 66.67, 55.83, 48.33, 2),
    ("DSS05", "Managed Sec. Services", 67.50, 51.67, 45.00, 2),
]

# Header
cols = [("Domain", 2.5), ("Level 1 (%)", 1.8), ("Level 2 (%)", 1.8), ("Level 3 (%)", 1.8), ("Cap. Level", 1.8)]
x = 1.0
for col_name, col_w in cols:
    box(s, x, 1.9, col_w-0.05, 0.5, C_ACCENT)
    txt(s, x, 1.92, col_w-0.05, 0.4, col_name, sz=12, color=C_BG, bold=True, align=PP_ALIGN.CENTER)
    x += col_w

for i, (code, name, l1, l2, l3, cap) in enumerate(results):
    y = 2.45 + i * 0.9
    bg = C_BG2 if i % 2 == 0 else C_BG
    # Domain
    box(s, 1.0, y, 2.45, 0.8, bg)
    txt(s, 1.1, y+0.05, 2.3, 0.3, code, sz=14, color=C_ACCENT, bold=True)
    txt(s, 1.1, y+0.35, 2.3, 0.3, name, sz=10, color=C_LIGHT)
    # Level 1
    box(s, 3.5, y, 1.75, 0.8, bg)
    l1_clr = C_YELLOW if l1 > 50 else C_ORANGE
    txt(s, 3.5, y+0.1, 1.75, 0.3, f"{l1:.1f}%", sz=16, color=l1_clr, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 3.5, y+0.45, 1.75, 0.3, "L ✓", sz=11, color=C_ACCENT2, align=PP_ALIGN.CENTER)
    # Level 2
    box(s, 5.3, y, 1.75, 0.8, bg)
    l2_clr = C_YELLOW if l2 > 50 else C_ORANGE
    txt(s, 5.3, y+0.1, 1.75, 0.3, f"{l2:.1f}%", sz=16, color=l2_clr, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 5.3, y+0.45, 1.75, 0.3, "L ✓", sz=11, color=C_ACCENT2, align=PP_ALIGN.CENTER)
    # Level 3
    box(s, 7.1, y, 1.75, 0.8, bg)
    txt(s, 7.1, y+0.1, 1.75, 0.3, f"{l3:.1f}%", sz=16, color=C_ORANGE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 7.1, y+0.45, 1.75, 0.3, "P ✗", sz=11, color=C_RED, align=PP_ALIGN.CENTER)
    # Cap Level
    box(s, 8.9, y, 1.75, 0.8, C_BG3)
    txt(s, 8.9, y+0.1, 1.75, 0.3, f"Level {cap}", sz=18, color=C_ACCENT, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 8.9, y+0.45, 1.75, 0.3, "Managed", sz=11, color=C_LIGHT, align=PP_ALIGN.CENTER)

txt(s, 1, 7.0, 10, 0.4, "→ Semua domain berada di Level 2: proses dasar sudah lengkap, tapi belum terstandarisasi formal.", sz=12, color=C_YELLOW)
snum(s, 10)


# ============================================================
# SLIDE 11: GAP ANALYSIS — VISUAL BAR CHART
# ============================================================
s = new_slide()
header(s, "05 | HASIL AUDIT", "GAP Analysis — As-is vs To-be")

domains_gap = ["APO12", "APO13", "BAI06", "DSS01", "DSS05"]
as_is = [2, 2, 2, 2, 2]
to_be = [4, 4, 4, 4, 4]

for i, (dom, cur, tgt) in enumerate(zip(domains_gap, as_is, to_be)):
    y = 2.0 + i * 1.0

    txt(s, 1, y+0.1, 1.5, 0.4, dom, sz=14, color=C_ACCENT, bold=True)

    # As-is bar (Blue)
    bar_w = cur * 1.8
    box(s, 2.5, y, bar_w, 0.35, C_ACCENT)
    txt(s, 2.5 + bar_w + 0.1, y, 1, 0.35, f"Level {cur}", sz=12, color=C_ACCENT, bold=True)

    # To-be bar (Teal, transparent effect)
    target_w = tgt * 1.8
    box(s, 2.5, y+0.4, target_w, 0.35, C_BG3)
    # Fill the target portion
    r_target = box(s, 2.5, y+0.4, target_w, 0.35, C_ACCENT2)
    r_target.fill.solid()
    r_target.fill.fore_color.rgb = C_ACCENT2
    # Add opacity effect by overlaying
    txt(s, 2.5 + target_w + 0.1, y+0.4, 1.5, 0.35, f"Level {tgt} (Target)", sz=12, color=C_ACCENT2, bold=True)

    # GAP badge
    gap = tgt - cur
    box(s, 11, y+0.1, 1.2, 0.6, C_RED)
    txt(s, 11, y+0.15, 1.2, 0.5, f"GAP\n{gap}", sz=14, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

# Legend
txt(s, 2.5, 7.0, 2, 0.3, "■ As-is (Saat Ini)", sz=11, color=C_ACCENT)
txt(s, 5.0, 7.0, 2, 0.3, "■ To-be (Target)", sz=11, color=C_ACCENT2)
txt(s, 7.5, 7.0, 3, 0.3, "■ GAP = 2 Level pada semua domain", sz=11, color=C_RED)
snum(s, 11)


# ============================================================
# SLIDE 12: TEMUAN UTAMA
# ============================================================
s = new_slide()
header(s, "05 | HASIL AUDIT", "Temuan Utama")

findings = [
    ("⚠️", "Dokumentasi Risiko Belum Lengkap",
     "Identifikasi risiko TI sudah dilakukan namun belum dikelola dalam risk register yang terstruktur",
     C_RED),
    ("🔒", "Sosialisasi Keamanan Terbatas",
     "Kebijakan keamanan informasi sudah ada namun program pelatihan keamanan belum berjalan rutin",
     C_ORANGE),
    ("🔄", "Change Management Belum Formal",
     "Pencatatan perubahan TI dilakukan, namun prosedur change management belum terdokumentasi formal",
     C_YELLOW),
    ("⚙️", "SOP Belum Konsisten",
     "Pemantauan operasional sudah berjalan rutin, namun SOP belum diterapkan konsisten di semua divisi",
     C_ACCENT),
    ("🛡️", "Audit Internal Belum Berkala",
     "Implementasi teknis keamanan sudah baik, namun audit keamanan internal belum dilaksanakan berkala",
     C_ACCENT2),
]
for i, (icon, title, desc, clr) in enumerate(findings):
    y = 1.9 + i * 1.05
    box(s, 1, y, 0.08, 0.85, clr)
    box(s, 1.1, y, 11.3, 0.85, C_BG2)
    txt(s, 1.3, y+0.05, 0.5, 0.4, icon, sz=20, color=C_WHITE)
    txt(s, 1.9, y+0.05, 10, 0.35, title, sz=14, color=clr, bold=True)
    txt(s, 1.9, y+0.4, 10, 0.35, desc, sz=12, color=C_LIGHT)
snum(s, 12)


# ============================================================
# SLIDE 13: REKOMENDASI — MENUJU LEVEL 3
# ============================================================
s = new_slide()
header(s, "06 | REKOMENDASI", "Menuju Level 3 (Established Process)")

txt(s, 1, 1.8, 11, 0.5,
    "Prioritas pertama: menyusun dokumentasi formal dan menstandarisasi prosedur tata kelola TI.",
    sz=14, color=C_LIGHT)

recs_l3 = [
    ("📋", "Dokumentasi Formal",
     "Menyusun kebijakan, prosedur, dan SOP tata kelola TI secara formal dan disahkan manajemen"),
    ("📊", "Risk Register",
     "Membuat risk register komprehensif mencakup risiko TI, dampak, probabilitas, dan mitigasi"),
    ("🎓", "Program Pelatihan",
     "Menyelenggarakan pelatihan keamanan informasi rutin untuk seluruh karyawan (min. 2x/tahun)"),
    ("🔄", "SOP Change Management",
     "Menyusun SOP manajemen perubahan TI formal mencakup alur persetujuan dan evaluasi"),
    ("📝", "Standarisasi Prosedur",
     "Menerapkan standar keamanan dan operasional yang konsisten di seluruh platform digital"),
]
for i, (icon, title, desc) in enumerate(recs_l3):
    y = 2.5 + i * 0.95
    box(s, 1, y, 11.5, 0.8, C_BG2)
    box(s, 1, y, 0.08, 0.8, C_ACCENT2)
    txt(s, 1.3, y+0.05, 0.5, 0.35, icon, sz=18, color=C_WHITE)
    txt(s, 1.9, y+0.05, 10, 0.35, title, sz=14, color=C_ACCENT2, bold=True)
    txt(s, 1.9, y+0.4, 10, 0.35, desc, sz=12, color=C_LIGHT)
snum(s, 13)


# ============================================================
# SLIDE 14: REKOMENDASI — MENUJU LEVEL 4
# ============================================================
s = new_slide()
header(s, "06 | REKOMENDASI", "Menuju Level 4 (Predictable Process)")

txt(s, 1, 1.8, 11, 0.5,
    "Setelah mencapai Level 3, targetkan pengukuran kinerja secara kuantitatif dan otomatisasi.",
    sz=14, color=C_LIGHT)

recs_l4 = [
    ("📈", "KPI Terukur",
     "Menetapkan Key Performance Indicators kuantitatif: uptime, MTTR, MTBF, jumlah insiden"),
    ("🖥️", "Monitoring Real-time",
     "Menerapkan dashboard monitoring otomatis untuk memantau kinerja platform digital"),
    ("🔍", "Audit Internal Berkala",
     "Menjadwalkan audit internal tata kelola TI secara periodik (minimal 2x per tahun)"),
    ("🧪", "Penetration Testing",
     "Melakukan penetration testing dan vulnerability assessment secara berkala"),
    ("📊", "SLA Internal",
     "Menetapkan Service Level Agreement internal untuk setiap layanan TI yang disediakan"),
]
for i, (icon, title, desc) in enumerate(recs_l4):
    y = 2.5 + i * 0.95
    box(s, 1, y, 11.5, 0.8, C_BG2)
    box(s, 1, y, 0.08, 0.8, C_ACCENT)
    txt(s, 1.3, y+0.05, 0.5, 0.35, icon, sz=18, color=C_WHITE)
    txt(s, 1.9, y+0.05, 10, 0.35, title, sz=14, color=C_ACCENT, bold=True)
    txt(s, 1.9, y+0.4, 10, 0.35, desc, sz=12, color=C_LIGHT)
snum(s, 14)


# ============================================================
# SLIDE 15: ROADMAP IMPLEMENTASI
# ============================================================
s = new_slide()
header(s, "06 | REKOMENDASI", "Roadmap Implementasi Bertahap")

# Timeline
phases = [
    ("FASE 1", "0-6 Bulan", "Menuju Level 3", C_ACCENT2, [
        "Dokumentasi kebijakan & SOP",
        "Pembentukan tim tata kelola TI",
        "Program pelatihan keamanan",
        "Standarisasi prosedur lintas divisi",
    ]),
    ("FASE 2", "6-12 Bulan", "Menuju Level 4", C_ACCENT, [
        "Implementasi KPI terukur",
        "Dashboard monitoring real-time",
        "Audit internal berkala",
        "Penetration testing terjadwal",
    ]),
]
for i, (phase, time, goal, clr, items) in enumerate(phases):
    l = 1 + i * 6.0
    # Phase header
    box(s, l, 1.9, 5.5, 0.7, clr)
    txt(s, l+0.2, 1.95, 1.5, 0.3, phase, sz=16, color=C_BG, bold=True)
    txt(s, l+0.2, 2.25, 1.5, 0.3, time, sz=12, color=C_BG)
    txt(s, l+2.5, 2.0, 3, 0.5, goal, sz=18, color=C_BG, bold=True, align=PP_ALIGN.RIGHT)
    # Items
    box(s, l, 2.7, 5.5, 3.5, C_BG2)
    for j, item in enumerate(items):
        y = 2.9 + j * 0.75
        box(s, l+0.2, y, 0.35, 0.35, clr)
        txt(s, l+0.2, y+0.02, 0.35, 0.3, "✓", sz=12, color=C_BG, bold=True, align=PP_ALIGN.CENTER)
        txt(s, l+0.7, y+0.02, 4.5, 0.3, item, sz=13, color=C_WHITE)

# Arrow between phases
txt(s, 6.2, 3.8, 0.8, 0.5, "→", sz=36, color=C_ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Target at bottom
box(s, 3, 6.5, 7, 0.6, C_BG3)
txt(s, 3, 6.55, 7, 0.5, "🎯  TARGET: Capability Level 4 (Predictable Process) dalam 12 bulan", sz=14, color=C_YELLOW, bold=True, align=PP_ALIGN.CENTER)
snum(s, 15)


# ============================================================
# SLIDE 16: KESIMPULAN & SARAN
# ============================================================
s = new_slide()
header(s, "07 | KESIMPULAN & SARAN", "Ringkasan Hasil Audit")

# Left: Kesimpulan
txt(s, 1, 1.8, 5.5, 0.4, "KESIMPULAN", sz=18, color=C_ACCENT2, bold=True)
aline(s, 1, 2.2, 2)
bullets(s, 1, 2.5, 5.5, 4, [
    "Seluruh 5 domain berada di Capability Level 2\n(Managed Process)",
    "GAP = 2 level menuju target Level 4\n(Predictable Process)",
    "DSS01 & DSS05 memiliki fondasi terbaik\n(Level 1 tertinggi: ~67%)",
    "APO12 memerlukan perhatian paling besar\n(Level 3 hanya 45.00%)",
], sz=12, color=C_LIGHT, spacing=12)

# Right: Saran
txt(s, 7, 1.8, 5.5, 0.4, "SARAN", sz=18, color=C_ACCENT, bold=True)
aline(s, 7, 2.2, 2)
bullets(s, 7, 2.5, 5.5, 4, [
    "Prioritas: dokumentasi formal semua\nkebijakan dan SOP tata kelola TI",
    "Terapkan KPI terukur untuk setiap\ndomain tata kelola TI",
    "Bentuk tim khusus tata kelola TI\ndan jalankan audit internal berkala",
    "Implementasi bertahap:\nFase 1 (6 bln) → Level 3\nFase 2 (12 bln) → Level 4",
], sz=12, color=C_LIGHT, spacing=12)
snum(s, 16)


# ============================================================
# SLIDE 17: TERIMA KASIH
# ============================================================
s = new_slide()
# Center content
txt(s, 1, 2.0, 11.5, 1.2, "Terima Kasih", sz=52, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
aline(s, 5.5, 3.3, 2.5)
txt(s, 1, 3.6, 11.5, 0.6, "Sesi Tanya Jawab", sz=26, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Team info
box(s, 3.5, 4.5, 6.5, 2.0, C_BG2)
txt(s, 3.5, 4.6, 6.5, 0.4, "TIM PENULIS", sz=14, color=C_ACCENT, bold=True, align=PP_ALIGN.CENTER)
txt(s, 3.5, 5.0, 6.5, 0.4, "Roki Anjas (11250066)  |  Susanto (11250068)", sz=14, color=C_WHITE, align=PP_ALIGN.CENTER)
txt(s, 3.5, 5.5, 6.5, 0.8,
    "Audit Sistem Informasi (453)\nJuarni Siregar, S.Pd., M.Kom\nUniversitas Nusa Mandiri | 2026",
    sz=12, color=C_LIGHT, align=PP_ALIGN.CENTER)
snum(s, 17)


# ============================================================
# SAVE
# ============================================================
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"✅ Presentasi PPTX berhasil disimpan: {OUTPUT}")
print(f"📊 Total slide: {TOTAL}")
print(f"🎨 Fitur:")
print(f"   - Dark navy theme dengan aksen biru")
print(f"   - 5 domain COBIT 2019 divisualisasikan")
print(f"   - GAP analysis dengan bar chart visual")
print(f"   - RACI Matrix interaktif")
print(f"   - Roadmap implementasi 2 fase")
print(f"   - Capability level per domain lengkap")
